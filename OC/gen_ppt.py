#!/usr/bin/env python3
"""
生成 OC Marketing 数字员工计划 PPT
逻辑线: 模型价格/路由 → 总方针 → 爆款收集 → 详细架构
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── 颜色系统 ──
BG_DARK    = RGBColor(0x1A, 0x1A, 0x2E)   # 深蓝黑背景
BG_CARD    = RGBColor(0x22, 0x22, 0x3A)   # 卡片背景
ACCENT     = RGBColor(0x6C, 0x63, 0xFF)   # 主紫色
ACCENT2    = RGBColor(0x00, 0xD2, 0xFF)   # 青色
GREEN      = RGBColor(0x00, 0xE6, 0x96)   # 绿色
ORANGE     = RGBColor(0xFF, 0x9F, 0x43)   # 橙色
RED        = RGBColor(0xFF, 0x6B, 0x6B)   # 红色
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
GRAY       = RGBColor(0xAA, 0xAA, 0xBB)
LIGHT_GRAY = RGBColor(0x88, 0x88, 0x99)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
W = prs.slide_width
H = prs.slide_height

# ══════════════════════════════════════════
# 工具函数
# ══════════════════════════════════════════

def add_bg(slide, color=BG_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape(slide, left, top, width, height, fill_color=None, line_color=None, line_width=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shp.fill.background()
    if fill_color:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill_color
    if line_color:
        shp.line.color.rgb = line_color
        shp.line.width = line_width or Pt(1)
    else:
        shp.line.fill.background()
    return shp

def add_rect(slide, left, top, width, height, fill_color=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shp.fill.background()
    if fill_color:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill_color
    shp.line.fill.background()
    return shp

def add_text(slide, left, top, width, height, text, font_size=18, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT):
    txbox = slide.shapes.add_textbox(left, top, width, height)
    tf = txbox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = alignment
    return txbox

def add_multiline(slide, left, top, width, height, lines, default_size=16, default_color=GRAY, line_spacing=1.3):
    """lines: list of (text, size, color, bold, alignment) or just str"""
    txbox = slide.shapes.add_textbox(left, top, width, height)
    tf = txbox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(lines):
        if isinstance(item, str):
            text, size, color, bold, align = item, default_size, default_color, False, PP_ALIGN.LEFT
        else:
            text = item[0]
            size = item[1] if len(item) > 1 else default_size
            color = item[2] if len(item) > 2 else default_color
            bold = item[3] if len(item) > 3 else False
            align = item[4] if len(item) > 4 else PP_ALIGN.LEFT
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.alignment = align
        p.space_after = Pt(size * (line_spacing - 1) * 2)
    return txbox

def add_card(slide, left, top, width, height, title, body_lines, accent=ACCENT, title_size=18):
    add_shape(slide, left, top, width, height, fill_color=BG_CARD, line_color=accent, line_width=Pt(1.5))
    # accent bar
    add_rect(slide, left, top, Pt(4), height, fill_color=accent)
    add_text(slide, left + Inches(0.25), top + Inches(0.12), width - Inches(0.4), Inches(0.4),
             title, font_size=title_size, color=accent, bold=True)
    add_multiline(slide, left + Inches(0.25), top + Inches(0.5), width - Inches(0.4), height - Inches(0.6),
                  body_lines, default_size=14, default_color=GRAY)

def add_arrow_right(slide, left, top, width=Inches(0.6), color=ACCENT):
    shp = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left, top, width, Inches(0.35))
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    return shp

def add_circle_number(slide, left, top, number, color=ACCENT):
    sz = Inches(0.55)
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, sz, sz)
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    tf = shp.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = str(number)
    p.font.size = Pt(20)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].space_before = Pt(0)
    return shp

def page_number(slide, num, total):
    add_text(slide, W - Inches(1.2), H - Inches(0.45), Inches(1), Inches(0.35),
             f"{num} / {total}", font_size=11, color=LIGHT_GRAY, alignment=PP_ALIGN.RIGHT)

TOTAL_SLIDES = 12

# ══════════════════════════════════════════
# SLIDE 1: 封面
# ══════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl)
# 装饰条
add_rect(sl, Inches(0), Inches(0), W, Inches(0.06), fill_color=ACCENT)
add_rect(sl, Inches(0), Inches(0.06), W, Inches(0.03), fill_color=ACCENT2)

add_text(sl, Inches(1), Inches(1.6), Inches(11), Inches(1),
         "基于 OpenClaw 的 Marketing 数字员工", font_size=42, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_text(sl, Inches(1), Inches(2.8), Inches(11), Inches(0.8),
         "模型选择 · 路由设计 · 爆款收集 · 系统架构", font_size=22, color=ACCENT2, alignment=PP_ALIGN.CENTER)

# 四个小标签
tags = ["模型定价", "智能路由", "爆款采集", "自动化架构"]
tag_colors = [ACCENT, GREEN, ORANGE, ACCENT2]
tag_start = Inches(3.2)
for i, (tag, tc) in enumerate(zip(tags, tag_colors)):
    x = tag_start + i * Inches(1.85)
    add_shape(sl, x, Inches(4.0), Inches(1.6), Inches(0.45), fill_color=None, line_color=tc, line_width=Pt(1.5))
    add_text(sl, x, Inches(4.02), Inches(1.6), Inches(0.45), tag, font_size=13, color=tc, bold=True, alignment=PP_ALIGN.CENTER)

add_text(sl, Inches(1), Inches(5.5), Inches(11), Inches(0.5),
         "2026-03 · Connact.ai 团队", font_size=14, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
page_number(sl, 1, TOTAL_SLIDES)

# ══════════════════════════════════════════
# SLIDE 2: 目录 / 逻辑线
# ══════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl)
add_text(sl, Inches(0.8), Inches(0.4), Inches(5), Inches(0.6), "演讲逻辑", font_size=30, color=WHITE, bold=True)

steps = [
    ("01", "模型价格与路由设计", "了解可用模型定价，设计三层路由策略", ACCENT),
    ("02", "Marketing 数字员工总方针", "产品理念、全流程工作流、OpenClaw 能力映射", GREEN),
    ("03", "P1 爆款收集 — 落地方案", "平台覆盖、采集流程、评分体系", ORANGE),
    ("04", "详细架构与实施计划", "系统架构、Skill 设计、成本监控、上线路径", ACCENT2),
]
for i, (num, title, desc, color) in enumerate(steps):
    y = Inches(1.5) + i * Inches(1.35)
    add_circle_number(sl, Inches(1.0), y, num, color)
    add_text(sl, Inches(1.8), y - Inches(0.02), Inches(4), Inches(0.45), title, font_size=22, color=WHITE, bold=True)
    add_text(sl, Inches(1.8), y + Inches(0.4), Inches(5), Inches(0.4), desc, font_size=14, color=GRAY)
    if i < 3:
        # 连接线（用小箭头暗示）
        add_text(sl, Inches(1.18), y + Inches(0.65), Inches(0.3), Inches(0.3), "│", font_size=14, color=color)

# 右侧示意
add_shape(sl, Inches(7.5), Inches(1.3), Inches(5), Inches(5.2), fill_color=BG_CARD, line_color=RGBColor(0x33,0x33,0x50))
add_multiline(sl, Inches(7.8), Inches(1.5), Inches(4.5), Inches(5),
    [
        ("📄 参考文档", 16, ACCENT2, True),
        (""),
        ("1. OpenClaw 可接入模型与定价研究报告", 14, WHITE, False),
        ("   → 22 页 · 模型定价/路由策略/配置教程", 12, GRAY),
        (""),
        ("2. P1 模型选择与路由方案（省钱版）", 14, WHITE, False),
        ("   → 成本估算/路由器代码/Prompt 模板", 12, GRAY),
        (""),
        ("3. P1 爆款收集执行计划", 14, WHITE, False),
        ("   → 采集流程/评分体系/Skill 设计/排期", 12, GRAY),
    ])
page_number(sl, 2, TOTAL_SLIDES)

# ══════════════════════════════════════════
# SLIDE 3: Part 1 标题页 — 模型价格与路由
# ══════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl)
add_rect(sl, Inches(0), Inches(0), Inches(0.15), H, fill_color=ACCENT)
add_text(sl, Inches(1), Inches(2.5), Inches(11), Inches(1),
         "Part 1", font_size=60, color=ACCENT, bold=True, alignment=PP_ALIGN.CENTER)
add_text(sl, Inches(1), Inches(3.8), Inches(11), Inches(0.8),
         "模型价格与路由设计", font_size=34, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_text(sl, Inches(2), Inches(4.8), Inches(9), Inches(0.6),
         "基于《OpenClaw 可接入模型与定价研究报告》的核心发现", font_size=16, color=GRAY, alignment=PP_ALIGN.CENTER)
page_number(sl, 3, TOTAL_SLIDES)

# ══════════════════════════════════════════
# SLIDE 4: 模型定价对比
# ══════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl)
add_text(sl, Inches(0.8), Inches(0.3), Inches(6), Inches(0.6), "模型定价三层架构", font_size=28, color=WHITE, bold=True)
add_text(sl, Inches(0.8), Inches(0.85), Inches(10), Inches(0.4),
         "把高价模型的调用压缩到真正需要的场景，其余用低价/免费模型覆盖", font_size=14, color=GRAY)

# 表头
header_y = Inches(1.5)
add_rect(sl, Inches(0.6), header_y, Inches(12), Inches(0.5), fill_color=RGBColor(0x2A, 0x2A, 0x45))
cols = [("层级", 1.3), ("模型", 3.0), ("输入$/MTok", 1.6), ("输出$/MTok", 1.6), ("上下文", 1.4), ("适用任务", 3.0)]
cx = Inches(0.7)
for label, w in cols:
    add_text(sl, cx, header_y + Inches(0.05), Inches(w), Inches(0.4), label, font_size=13, color=ACCENT2, bold=True)
    cx += Inches(w)

# 数据行
rows = [
    ("🟢 轻量", "Z.AI GLM-4.7-Flash", "免费", "免费", "128K", "日报/简单文本", GREEN),
    ("🟢 轻量", "Z.AI GLM-4.7-FlashX", "$0.07", "$0.40", "200K", "快速分类/标签", GREEN),
    ("🟢 轻量", "OpenAI gpt-5-nano", "$0.05", "$0.40", "128K", "轻量备选", GREEN),
    ("🟡 中等", "OpenAI gpt-5-mini", "$0.25", "$2.00", "128K", "深度分析/常规编码", ORANGE),
    ("🟡 中等", "MiniMax M2.5", "$0.30", "$1.20", "200K", "中文深度分析", ORANGE),
    ("🔴 高级", "Anthropic Sonnet 4.6", "$3.00", "$15.00", "200K", "月度复盘/复杂Agent", RED),
    ("🔴 高级", "OpenAI gpt-5.3-codex", "$1.75", "$14.00", "400K", "大仓库/系统设计", RED),
]
for i, (tier, model, inp, out, ctx, task, color) in enumerate(rows):
    ry = header_y + Inches(0.55) + i * Inches(0.52)
    if i % 2 == 0:
        add_rect(sl, Inches(0.6), ry, Inches(12), Inches(0.5), fill_color=RGBColor(0x20, 0x20, 0x35))
    vals = [tier, model, inp, out, ctx, task]
    cx = Inches(0.7)
    for j, (_, w) in enumerate(cols):
        c = color if j == 0 else (GREEN if vals[j] == "免费" else GRAY)
        add_text(sl, cx, ry + Inches(0.05), Inches(w), Inches(0.4), vals[j], font_size=12, color=c)
        cx += Inches(w)

# 关键洞察
add_shape(sl, Inches(0.6), Inches(5.8), Inches(12), Inches(1.2), fill_color=BG_CARD, line_color=ACCENT)
add_multiline(sl, Inches(0.9), Inches(5.9), Inches(11.5), Inches(1.1), [
    ("💡 关键洞察", 16, ACCENT, True),
    ("• 轻量层 vs 高级层价格差距可达 70-350 倍 — 路由策略的意义巨大", 13, WHITE),
    ("• Z.AI Flash 免费档是日报生成的绝佳选择；FlashX 的 $0.07 是分类任务的性价比之王", 13, WHITE),
    ("• OpenRouter 可做多 Provider 冗余（5.5% 手续费），但 P1 阶段量小，直连更省", 13, WHITE),
])
page_number(sl, 4, TOTAL_SLIDES)

# ══════════════════════════════════════════
# SLIDE 5: 路由策略设计
# ══════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl)
add_text(sl, Inches(0.8), Inches(0.3), Inches(8), Inches(0.6), "静态 + 动态路由策略", font_size=28, color=WHITE, bold=True)

# 左: 静态路由
add_card(sl, Inches(0.6), Inches(1.2), Inches(5.8), Inches(5.5),
    "🔧 静态路由（规则驱动）", [
        ("规则 1 → 图片输入 → 高级层（支持多模态）", 14, WHITE),
        ("规则 2 → 短文本 < 800字 + 无复杂关键词 → 轻量层", 14, WHITE),
        ("规则 3 → 编码/修复/写测试指令 → 中等层", 14, WHITE),
        ("规则 4 → 全仓库/架构/安全审计 → 高级层", 14, WHITE),
        (""),
        ("✅ 优点: 可解释 · 可审计 · 零额外成本", 13, GREEN),
        ("⚠️  建议: 路由做到「会话/作业级别」", 13, ORANGE),
        ("　避免一个 session 内频繁换模型导致上下文漂移", 12, GRAY),
    ], accent=ACCENT)

# 右: 动态路由
add_card(sl, Inches(6.9), Inches(1.2), Inches(5.8), Inches(5.5),
    "⚡ 动态路由（指标驱动）", [
        ("实时指标采集:", 14, ACCENT2, True),
        ("• 延迟: TTFT / p50 / p95 总耗时", 13, WHITE),
        ("• 错误: 429/5xx 比例 · 超时 · billing 错误", 13, WHITE),
        ("• 成本: 日/月累计 · 预算剩余百分比", 13, WHITE),
        ("• 配额: 订阅制剩余请求 · 距刷新时间", 13, WHITE),
        (""),
        ("动态决策规则:", 14, ACCENT2, True),
        ("• 预算剩余 < 15% → 高级→中等 降级", 13, ORANGE),
        ("• 预算剩余 < 5% → 中等→轻量 降级", 13, ORANGE),
        ("• p95 > 10s(轻/中) 或 > 25s(高) → 熔断 10min", 13, RED),
        ("• 429 连续 ≥3 次/2min → 熔断 + 切 Provider", 13, RED),
    ], accent=ACCENT2)
page_number(sl, 5, TOTAL_SLIDES)

# ══════════════════════════════════════════
# SLIDE 6: Part 2 标题页 — Marketing 数字员工总方针
# ══════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl)
add_rect(sl, Inches(0), Inches(0), Inches(0.15), H, fill_color=GREEN)
add_text(sl, Inches(1), Inches(2.5), Inches(11), Inches(1),
         "Part 2", font_size=60, color=GREEN, bold=True, alignment=PP_ALIGN.CENTER)
add_text(sl, Inches(1), Inches(3.8), Inches(11), Inches(0.8),
         "Marketing 数字员工总方针", font_size=34, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_text(sl, Inches(2), Inches(4.8), Inches(9), Inches(0.6),
         "OpenClaw 驱动的内容生产自动化流水线", font_size=16, color=GRAY, alignment=PP_ALIGN.CENTER)
page_number(sl, 6, TOTAL_SLIDES)

# ══════════════════════════════════════════
# SLIDE 7: 全流程工作流
# ══════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl)
add_text(sl, Inches(0.8), Inches(0.3), Inches(8), Inches(0.6), "内容生产全流程（9 阶段）", font_size=28, color=WHITE, bold=True)

# 上排5个
flow_top = [
    ("P1", "爆款收集", "自动抓取\n多平台热门内容", ACCENT, True),
    ("P2", "内容筛选", "周会评审\n团队投票决策", GREEN, False),
    ("P3", "脚本改编", "LLM 辅助\n改编为品牌脚本", ORANGE, False),
    ("P4", "拍摄制作", "按脚本执行\n拍摄+录制", LIGHT_GRAY, False),
    ("P5", "剪辑优化", "AI 字幕/配乐\n节奏优化", LIGHT_GRAY, False),
]
# 下排4个
flow_bot = [
    ("P6", "发布上线", "多平台定时发布", LIGHT_GRAY, False),
    ("P7", "数据监测", "实时追踪\n播放/互动数据", LIGHT_GRAY, False),
    ("P8", "周复盘", "数据驱动\n总结成功模式", LIGHT_GRAY, False),
    ("P9", "优化迭代", "反馈闭环\n持续提升命中率", LIGHT_GRAY, False),
]

for i, (num, title, desc, color, highlight) in enumerate(flow_top):
    x = Inches(0.5) + i * Inches(2.5)
    y = Inches(1.4)
    bw = Inches(2.1)
    bh = Inches(1.9)
    if highlight:
        add_shape(sl, x, y, bw, bh, fill_color=RGBColor(0x2A, 0x20, 0x55), line_color=ACCENT, line_width=Pt(2.5))
        add_text(sl, x + Inches(0.6), y + Inches(0.1), Inches(1.5), Inches(0.35), "★ 当前重点", font_size=10, color=ACCENT, bold=True)
    else:
        add_shape(sl, x, y, bw, bh, fill_color=BG_CARD, line_color=RGBColor(0x44,0x44,0x55))
    add_text(sl, x + Inches(0.12), y + Inches(0.1), Inches(0.5), Inches(0.35), num, font_size=14, color=color, bold=True)
    add_text(sl, x + Inches(0.12), y + Inches(0.5), bw - Inches(0.2), Inches(0.4), title, font_size=18, color=WHITE, bold=True)
    add_text(sl, x + Inches(0.12), y + Inches(1.0), bw - Inches(0.2), Inches(0.8), desc, font_size=12, color=GRAY)
    if i < 4:
        add_text(sl, x + bw + Inches(0.05), y + Inches(0.7), Inches(0.4), Inches(0.4), "→", font_size=18, color=LIGHT_GRAY)

# 连接箭头: P5 → P6
add_text(sl, Inches(11.3), Inches(3.6), Inches(0.8), Inches(0.8), "↓", font_size=24, color=LIGHT_GRAY)

for i, (num, title, desc, color, _) in enumerate(flow_bot):
    x = Inches(0.5) + (3 - i) * Inches(2.5) + Inches(2.5)  # 右到左排列以形成 U 形
    # 实际还是左到右更清晰
    x = Inches(0.5) + i * Inches(2.5) + Inches(2.5)
    y = Inches(4.3)
    bw = Inches(2.1)
    bh = Inches(1.5)
    add_shape(sl, x, y, bw, bh, fill_color=BG_CARD, line_color=RGBColor(0x44,0x44,0x55))
    add_text(sl, x + Inches(0.12), y + Inches(0.1), Inches(0.5), Inches(0.35), num, font_size=14, color=color, bold=True)
    add_text(sl, x + Inches(0.12), y + Inches(0.45), bw - Inches(0.2), Inches(0.4), title, font_size=18, color=WHITE, bold=True)
    add_text(sl, x + Inches(0.12), y + Inches(0.9), bw - Inches(0.2), Inches(0.6), desc, font_size=12, color=GRAY)
    if i < 3:
        add_text(sl, x + bw + Inches(0.05), y + Inches(0.5), Inches(0.4), Inches(0.4), "→", font_size=18, color=LIGHT_GRAY)

# OpenClaw 能力映射
add_shape(sl, Inches(0.5), Inches(6.2), Inches(12.3), Inches(0.95), fill_color=BG_CARD, line_color=GREEN)
add_multiline(sl, Inches(0.8), Inches(6.3), Inches(11.8), Inches(0.8), [
    ("OpenClaw 能力映射:  ", 14, GREEN, True),
    ("Browser Control → P1 抓取 ｜ Cron → P1/P6/P7 定时 ｜ Memory → P1/P8/P9 学习闭环 ｜ Skills → 全流程插件化 ｜ Chat → P2 周会/日报推送", 12, WHITE),
])
page_number(sl, 7, TOTAL_SLIDES)

# ══════════════════════════════════════════
# SLIDE 8: Part 3 标题页 — P1 爆款收集
# ══════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl)
add_rect(sl, Inches(0), Inches(0), Inches(0.15), H, fill_color=ORANGE)
add_text(sl, Inches(1), Inches(2.5), Inches(11), Inches(1),
         "Part 3", font_size=60, color=ORANGE, bold=True, alignment=PP_ALIGN.CENTER)
add_text(sl, Inches(1), Inches(3.8), Inches(11), Inches(0.8),
         "P1 爆款收集 — 落地方案", font_size=34, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_text(sl, Inches(2), Inches(4.8), Inches(9), Inches(0.6),
         "全自动 · 多平台 · 可学习的爆款内容采集系统", font_size=16, color=GRAY, alignment=PP_ALIGN.CENTER)
page_number(sl, 8, TOTAL_SLIDES)

# ══════════════════════════════════════════
# SLIDE 9: P1 采集流程 + 两阶段 LLM 路由
# ══════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl)
add_text(sl, Inches(0.8), Inches(0.3), Inches(8), Inches(0.6), "P1 采集流程 × 两阶段 LLM 路由", font_size=28, color=WHITE, bold=True)

# 左侧: 9 步流程
steps_data = [
    ("①", "关键词加载", "无 LLM", GRAY),
    ("②", "平台抓取 (Browser)", "无 LLM", GRAY),
    ("③", "数据抽取 → JSON", "无 LLM", GRAY),
    ("④", "去重 (Memory)", "无 LLM", GRAY),
    ("⑤", "硬指标打分", "无 LLM", GRAY),
    ("⑥", "软指标分析", "LLM ✦", ACCENT),
    ("⑦", "综合评分 A/B/C", "无 LLM", GRAY),
    ("⑧", "存储 + Memory", "无 LLM", GRAY),
    ("⑨", "日报推送", "LLM ✦", ACCENT),
]
for i, (num, name, llm, color) in enumerate(steps_data):
    y = Inches(1.15) + i * Inches(0.63)
    # 步骤条
    bar_color = RGBColor(0x2A, 0x20, 0x55) if color == ACCENT else RGBColor(0x22, 0x22, 0x35)
    add_rect(sl, Inches(0.5), y, Inches(4.3), Inches(0.5), fill_color=bar_color)
    add_text(sl, Inches(0.6), y + Inches(0.05), Inches(0.4), Inches(0.4), num, font_size=14, color=color, bold=True)
    add_text(sl, Inches(1.1), y + Inches(0.05), Inches(2.5), Inches(0.4), name, font_size=14, color=WHITE)
    add_text(sl, Inches(3.5), y + Inches(0.05), Inches(1.2), Inches(0.4), llm, font_size=12, color=color, bold=color==ACCENT)

# 右侧: 两阶段漏斗
add_card(sl, Inches(5.3), Inches(1.1), Inches(7.4), Inches(5.8),
    "⑥ 两阶段 LLM 漏斗（核心省钱策略）", [
        ("", 6, GRAY),
        ("阶段 1 — 快速分类（全量）", 16, GREEN, True),
        ("模型: GLM-4.7-FlashX（$0.07/MTok）", 14, WHITE),
        ("任务: Hook 类型 / 情绪标签 / 品牌相关度", 13, GRAY),
        ("成本: ¥0.003/条 · 30条/天 = ¥0.09/天", 13, GREEN),
        ("输出: JSON → {hook, emotion, brand_rel, worth_deep}", 12, GRAY),
        (""),
        ("　　　　　  ⬇  brand_relevance ≥ 0.5 才进入 (~40%)", 13, ORANGE),
        (""),
        ("阶段 2 — 深度分析（条件触发）", 16, ACCENT2, True),
        ("模型: gpt-5-mini（$0.25/$2.00/MTok）", 14, WHITE),
        ("任务: 内容公式 / Hook 细节 / 可复制性 / 改编建议", 13, GRAY),
        ("成本: ¥0.02/条 · 12条/天 = ¥0.24/天", 13, ACCENT2),
        (""),
        ("⑨ 日报趋势总结", 16, GREEN, True),
        ("模型: GLM-4.7-Flash（免费！）", 14, GREEN),
        ("成本: ¥0/天", 13, GREEN),
    ], accent=ORANGE, title_size=16)
page_number(sl, 9, TOTAL_SLIDES)

# ══════════════════════════════════════════
# SLIDE 10: 成本估算
# ══════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl)
add_text(sl, Inches(0.8), Inches(0.3), Inches(8), Inches(0.6), "月度成本估算 — 三档方案", font_size=28, color=WHITE, bold=True)

# 三个方案卡片
plans = [
    ("🟢 极省版", "≈ ¥3/月", [
        "全量用 Z.AI Flash 免费档",
        "FlashX 做简单分类",
        "不做深度分析",
        "",
        "适合: MVP 验证阶段",
    ], GREEN),
    ("🟡 推荐版", "≈ ¥10/月", [
        "FlashX 快速分类 → ¥3",
        "gpt-5-mini 深度分析 → ¥7",
        "Flash 免费做日报 → ¥0",
        "",
        "适合: 正式运营（推荐）",
    ], ORANGE),
    ("🔴 高配版", "≈ ¥32/月", [
        "gpt-5-mini 全量分析",
        "Sonnet 4.6 月度复盘",
        "追求最高分析质量",
        "",
        "适合: 预算充裕团队",
    ], RED),
]
for i, (title, price, details, color) in enumerate(plans):
    x = Inches(0.5) + i * Inches(4.2)
    w = Inches(3.9)
    add_shape(sl, x, Inches(1.2), w, Inches(4.0), fill_color=BG_CARD, line_color=color, line_width=Pt(2))
    add_text(sl, x + Inches(0.2), Inches(1.35), w - Inches(0.4), Inches(0.4), title, font_size=20, color=color, bold=True)
    add_text(sl, x + Inches(0.2), Inches(1.85), w - Inches(0.4), Inches(0.6), price, font_size=36, color=WHITE, bold=True)
    lines = [(d, 14, GRAY if d else GRAY) for d in details]
    add_multiline(sl, x + Inches(0.2), Inches(2.7), w - Inches(0.4), Inches(2.3), lines)
    if i == 1:  # 推荐标记
        add_shape(sl, x + w - Inches(1.3), Inches(1.2), Inches(1.3), Inches(0.35), fill_color=ORANGE)
        add_text(sl, x + w - Inches(1.3), Inches(1.2), Inches(1.3), Inches(0.35), "推荐", font_size=13, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# 对比条
add_shape(sl, Inches(0.5), Inches(5.5), Inches(12.3), Inches(1.5), fill_color=BG_CARD, line_color=ACCENT)
add_multiline(sl, Inches(0.8), Inches(5.65), Inches(11.8), Inches(1.3), [
    ("📊 成本对比（假设 1500 条/月分析量）", 16, ACCENT, True),
    (""),
    ("┃ 全量 Opus 4.6:  ¥375/月   ┃ 全量 gpt-5-mini:  ¥30/月   ┃ 推荐分层路由:  ¥10/月   ┃ 极省版:  ¥3/月 ┃", 14, WHITE),
    ("　 节省 97%+                     节省 67%                       最佳性价比 ★                 MVP 验证", 12, GREEN),
])
page_number(sl, 10, TOTAL_SLIDES)

# ══════════════════════════════════════════
# SLIDE 11: Part 4 标题页 — 详细架构
# ══════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl)
add_rect(sl, Inches(0), Inches(0), Inches(0.15), H, fill_color=ACCENT2)
add_text(sl, Inches(1), Inches(2.5), Inches(11), Inches(1),
         "Part 4", font_size=60, color=ACCENT2, bold=True, alignment=PP_ALIGN.CENTER)
add_text(sl, Inches(1), Inches(3.8), Inches(11), Inches(0.8),
         "详细架构与实施计划", font_size=34, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_text(sl, Inches(2), Inches(4.8), Inches(9), Inches(0.6),
         "系统架构 · Skill 设计 · 成本监控 · 上线路径", font_size=16, color=GRAY, alignment=PP_ALIGN.CENTER)
page_number(sl, 11, TOTAL_SLIDES)

# ══════════════════════════════════════════
# SLIDE 12: 系统架构 + 实施排期
# ══════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl)
add_text(sl, Inches(0.8), Inches(0.3), Inches(8), Inches(0.6), "系统架构 + 6 周实施路径", font_size=28, color=WHITE, bold=True)

# 左: 架构图（用卡片模拟）
arch_x = Inches(0.4)
arch_w = Inches(6.5)

# 用户入口层
add_shape(sl, arch_x, Inches(1.1), arch_w, Inches(0.7), fill_color=RGBColor(0x1E, 0x3A, 0x5F), line_color=ACCENT2)
add_text(sl, arch_x + Inches(0.15), Inches(1.2), arch_w, Inches(0.5),
         "📱 用户入口 — Telegram / 飞书 / Discord → OpenClaw Gateway", font_size=13, color=ACCENT2, bold=True)

# 路由层
add_shape(sl, arch_x, Inches(2.0), arch_w, Inches(0.9), fill_color=RGBColor(0x2A, 0x1E, 0x5F), line_color=ACCENT)
add_multiline(sl, arch_x + Inches(0.15), Inches(2.05), arch_w, Inches(0.8), [
    ("🔀 路由与策略层", 14, ACCENT, True),
    ("Router Service（静态+动态）  ←→  路由配置 router.yaml  ←→  指标存储（延迟/错误/成本）", 11, GRAY),
])

# Skills 执行层
add_shape(sl, arch_x, Inches(3.1), arch_w, Inches(1.6), fill_color=RGBColor(0x1E, 0x3A, 0x2A), line_color=GREEN)
add_multiline(sl, arch_x + Inches(0.15), Inches(3.15), arch_w, Inches(1.5), [
    ("⚙️  OpenClaw Skills 执行层", 14, GREEN, True),
    (""),
    ("viral-collector          viral-analyzer           viral-reporter",  12, WHITE),
    ("Cron 每6h · Browser     两阶段 LLM 漏斗          模板渲染 + LLM 趋势",  11, GRAY),
    ("多平台抓取 · 去重       FlashX→gpt-5-mini        Flash(免费) · 推送", 11, GRAY),
])

# 存储层
add_shape(sl, arch_x, Inches(4.9), Inches(3.1), Inches(0.9), fill_color=BG_CARD, line_color=ORANGE)
add_multiline(sl, arch_x + Inches(0.15), Inches(4.95), Inches(2.9), Inches(0.8), [
    ("💾 数据存储", 13, ORANGE, True),
    ("/data/oc/viral_pool/", 11, GRAY),
    ("raw/ · analyzed/ · reports/", 11, GRAY),
])

add_shape(sl, arch_x + Inches(3.3), Inches(4.9), Inches(3.2), Inches(0.9), fill_color=BG_CARD, line_color=ACCENT2)
add_multiline(sl, arch_x + Inches(3.45), Inches(4.95), Inches(3.0), Inches(0.8), [
    ("🧠 Persistent Memory", 13, ACCENT2, True),
    ("viral_urls_seen · patterns", 11, GRAY),
    ("preferences · keywords", 11, GRAY),
])

# 模型层
add_shape(sl, arch_x, Inches(6.0), arch_w, Inches(0.8), fill_color=RGBColor(0x3A, 0x1E, 0x1E), line_color=RED)
add_multiline(sl, arch_x + Inches(0.15), Inches(6.05), arch_w, Inches(0.7), [
    ("☁️  模型 API 层", 13, RED, True),
    ("Z.AI Flash/FlashX(免费/轻量)   |   OpenAI gpt-5-mini(中等)   |   Anthropic Sonnet(高级·可选)", 11, GRAY),
])

# 连接箭头
for y_pos in [Inches(1.82), Inches(2.92), Inches(4.72), Inches(5.82)]:
    add_text(sl, arch_x + Inches(2.8), y_pos, Inches(1), Inches(0.2), "▼", font_size=14, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# 右: 实施排期
rx = Inches(7.3)
rw = Inches(5.6)
add_text(sl, rx, Inches(1.0), rw, Inches(0.4), "📅 6 周实施路径", font_size=18, color=WHITE, bold=True)

phases = [
    ("Phase 1 — MVP", "第 1-2 周", [
        "安装配置 OpenClaw",
        "编写 viral-collector（抖音+小红书）",
        "确定关键词 + 爆款阈值 V1",
        "手动触发采集 · 验证数据质量",
    ], ACCENT),
    ("Phase 2 — 自动化", "第 3-4 周", [
        "编写 viral-analyzer + reporter",
        "配置 Cron 定时任务",
        "接入 Telegram/飞书推送",
        "扩展到 B站 + YouTube",
    ], GREEN),
    ("Phase 3 — 学习闭环", "第 5-6 周", [
        "P2 筛选结果回流 Memory",
        "关键词效果分析 + 权重自调整",
        "阈值动态调整 · 模式总结",
        "月度套路手册输出",
    ], ORANGE),
]

for i, (title, time_range, items, color) in enumerate(phases):
    y = Inches(1.5) + i * Inches(1.85)
    add_shape(sl, rx, y, rw, Inches(1.65), fill_color=BG_CARD, line_color=color)
    add_rect(sl, rx, y, Pt(4), Inches(1.65), fill_color=color)
    add_text(sl, rx + Inches(0.2), y + Inches(0.08), Inches(3), Inches(0.3), title, font_size=15, color=color, bold=True)
    add_text(sl, rx + rw - Inches(1.4), y + Inches(0.08), Inches(1.2), Inches(0.3), time_range, font_size=12, color=GRAY, alignment=PP_ALIGN.RIGHT)
    lines = [(f"• {item}", 12, WHITE) for item in items]
    add_multiline(sl, rx + Inches(0.2), y + Inches(0.42), rw - Inches(0.4), Inches(1.15), lines, line_spacing=1.2)

page_number(sl, 12, TOTAL_SLIDES)

# ══════════════════════════════════════════
# 保存
# ══════════════════════════════════════════
output_path = "/Users/linlin/Downloads/Coldemail-agent-dev/OC/Marketing数字员工计划.pptx"
prs.save(output_path)
print(f"✅ PPT saved to: {output_path}")
print(f"   Total slides: {len(prs.slides)}")
